from pptx import Presentation # type: ignore
from pptx.util import Pt # type: ignore

# Function to read text from a file
def read_text_from_file(file_path):
    with open(file_path, "r", encoding="utf-8") as file:
        return file.read().strip()

# Function to create PowerPoint from text
def create_ppt_from_text(file_path, output_file):
    slide_data = read_text_from_file(file_path)
    prs = Presentation()

    slides = slide_data.split("\n\n")  # Splitting slides by double newlines
    
    for i, slide in enumerate(slides):
        lines = [line.strip() for line in slide.split("\n") if line.strip()]  # Remove empty lines

        if not lines:
            continue  # Skip empty sections

        # Ensure the first line is a title
        if lines[0].startswith("Slide Title:"):
            slide_layout = prs.slide_layouts[1]  # Title & Content layout
            slide_obj = prs.slides.add_slide(slide_layout)
            
            # Extract slide title
            slide_title = lines[0].replace("Slide Title:", "").strip()
            slide_obj.shapes.title.text = slide_title

            # Adjust font size for title
            title_shape = slide_obj.shapes.title
            title_shape.text_frame.text = slide_title
            title_shape.text_frame.paragraphs[0].font.size = Pt(36 if i == 0 else 28)  # Reduce font size for the first slide

            # Extract bullet points
            content_placeholder = slide_obj.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()  # Remove default text

            # Start with the first bullet point
            first_bullet_added = False

            for line in lines[1:]:
                if line.startswith("-"):  # Checking for valid bullet points
                    bullet_point = line.lstrip("- ").strip()  # Remove "- " and extra spaces
                    if not first_bullet_added:
                        text_frame.text = bullet_point  # First bullet point
                        text_frame.paragraphs[0].font.size = Pt(24 if i == 0 else 20)  # Adjust font size
                        first_bullet_added = True
                    else:
                        p = text_frame.add_paragraph()
                        p.text = bullet_point  # Add subsequent bullet points
                        p.level = 0  # Ensure it's a primary bullet point
                        p.font.size = Pt(24 if i == 0 else 20)  # Adjust bullet point font size

    if len(prs.slides) == 0:
        print("⚠️ No slides were created. Please check the input format!")

    prs.save(output_file)
    print(f"🎉 PowerPoint file successfully created: {output_file}")

# Example usage
create_ppt_from_text("slides.txt", "output.pptx")
