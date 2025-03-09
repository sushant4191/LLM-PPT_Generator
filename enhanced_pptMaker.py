from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_ppt_from_text(slides_file, output_ppt):
    # Create a PowerPoint presentation
    prs = Presentation()

    # Define slide layout (Title + Content)
    slide_layout = prs.slide_layouts[1]

    # Read the slides content
    with open(slides_file, "r", encoding="utf-8") as file:
        content = file.read().strip()

    # Split content into slides
    slides = content.split("\n\n")

    for slide_text in slides:
        lines = slide_text.strip().split("\n")

        if len(lines) < 2:
            continue  # Skip malformed slides

        # Extract title
        title_text = lines[0].replace("Slide Title: ", "").strip()
        bullet_points = [line.replace("- ", "").strip() for line in lines[1:] if "Bullet Points:" not in line]  # Remove dashes and filter "Bullet Points:"

        # Create slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray background

        # Set title properties
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Clear any existing text in the content placeholder
        content.text = ""
        for point in bullet_points:
            p = content.text_frame.add_paragraph()
            p.text = point
            p.space_after = Pt(10)
            p.level = 0
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(50, 50, 50)  # Dark gray

    prs.save(output_ppt)
    print(f"✅ PowerPoint saved as {output_ppt}")

create_ppt_from_text("slides.txt", "output.pptx")
